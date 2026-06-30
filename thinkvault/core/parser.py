"""
文档解析器 — PDF / Word / PPT / Excel / TXT / Markdown / 音频 / 视频
"""

import os
import threading
import time
import shutil
import tempfile
from pathlib import Path
from typing import Optional
from dataclasses import dataclass, field

from thinkvault.utils.logger import logger

# DOCX 锁等待超时（秒），通过环境变量 THINKVAULT_DOCX_TIMEOUT 配置
DOCX_LOCK_TIMEOUT = int(os.environ.get("THINKVAULT_DOCX_TIMEOUT", "5"))

# ── 可选依赖：faster-whisper ──
try:
    from faster_whisper import WhisperModel
    _whisper_available = True
except ImportError:
    _whisper_available = False

# Whisper 模型缓存（避免每次音频解析重新加载）
_whisper_model_cache: dict = {}
_whisper_lock = threading.Lock()

# ── 可选依赖：ffmpeg ──
_ffmpeg_available = shutil.which("ffmpeg") is not None


@dataclass
class ParsedDocument:
    """解析后的文档结构"""
    file_path: str
    file_name: str
    file_type: str       # pdf / docx / txt / md / mp3 / mp4 ...
    total_pages: int = 0
    paragraphs: list[dict] = field(default_factory=list)
    tables: list[dict] = field(default_factory=list)
    raw_text: str = ""
    parse_error: Optional[str] = None

    @property
    def is_empty(self) -> bool:
        return len(self.raw_text.strip()) == 0


# 音频文件扩展名
AUDIO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".flac", ".ogg"}
# 视频文件扩展名
VIDEO_EXTENSIONS = {".mp4", ".mkv", ".avi", ".mov", ".webm"}


class DocumentParser:
    """多格式文档解析器"""

    SUPPORTED_TYPES = {
        # 文档类
        ".pdf", ".docx", ".pptx", ".xlsx", ".xlsm", ".txt", ".md", ".markdown",
        # 音频类
        ".mp3", ".wav", ".m4a", ".flac", ".ogg",
        # 视频类
        ".mp4", ".mkv", ".avi", ".mov", ".webm",
    }

    @classmethod
    def parse(cls, file_path: str) -> ParsedDocument:
        path = Path(file_path)
        if not path.exists():
            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type="",
                parse_error=f"文件不存在: {file_path}",
            )

        suffix = path.suffix.lower()
        if suffix not in cls.SUPPORTED_TYPES:
            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type=suffix,
                parse_error=f"不支持的文件格式: {suffix}，目前支持 {cls.SUPPORTED_TYPES}",
            )

        if suffix == ".pdf":
            result = cls._parse_pdf(file_path, path)
        elif suffix == ".docx":
            result = cls._parse_docx(file_path, path)
        elif suffix == ".pptx":
            result = cls._parse_pptx(file_path, path)
        elif suffix in (".xlsx", ".xlsm"):
            result = cls._parse_xlsx(file_path, path)
        elif suffix in (".txt",):
            result = cls._parse_txt(file_path, path)
        elif suffix in (".md", ".markdown"):
            result = cls._parse_markdown(file_path, path)
        elif suffix in AUDIO_EXTENSIONS:
            result = cls._parse_audio(file_path, path)
        elif suffix in VIDEO_EXTENSIONS:
            result = cls._parse_video(file_path, path)
        else:
            result = ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type=suffix,
                parse_error=f"未处理的文件格式: {suffix}",
            )

        # MarkItDown 补充解析：根据模式决定是否用 MarkItDown 替换/兜底
        from thinkvault.core import markitdown_adapter
        return markitdown_adapter.convert_with_fallback(file_path, result)

    @classmethod
    def _parse_pdf(cls, file_path: str, path: Path) -> ParsedDocument:
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(file_path)
            try:
                paragraphs = []
                full_text_parts = []
                table_count = 0
                # 记录每页 PyMuPDF 提取到的文本长度，用于混合型 PDF 去重
                page_text_lengths: dict[int, int] = {}

                for page_num, page in enumerate(doc):
                    page_char_count = 0
                    # 提取文本块（保留段落结构）
                    blocks = page.get_text("dict")["blocks"]
                    for block in blocks:
                        if block["type"] == 0:  # 文本块
                            text = ""
                            for line in block.get("lines", []):
                                spans_text = "".join(
                                    span["text"] for span in line.get("spans", [])
                                )
                                text += spans_text
                            text = text.strip()
                            if text:
                                paragraphs.append({
                                    "page": page_num + 1,
                                    "text": text,
                                    "char_count": len(text),
                                })
                                full_text_parts.append(text)
                                page_char_count += len(text)
                        elif block["type"] == 1:  # 图片块
                            pass

                    # 提取表格
                    tables_on_page = page.find_tables()
                    if tables_on_page:
                        for tab in tables_on_page:
                            table_count += 1
                            table_data = tab.extract()
                            if table_data:
                                header = table_data[0]
                                rows = table_data[1:]
                                lines = []
                                if header:
                                    lines.append(" | ".join(str(c) for c in header))
                                    lines.append(" | ".join("---" for _ in header))
                                for row in rows:
                                    lines.append(" | ".join(str(c) for c in row))
                                table_md = "\n".join(lines)
                                paragraphs.append({
                                    "page": page_num + 1,
                                    "text": table_md,
                                    "char_count": len(table_md),
                                    "is_table": True,
                                })
                                full_text_parts.append(table_md)
                                page_char_count += len(table_md)

                    page_text_lengths[page_num] = page_char_count

                total_pages = len(doc)
                raw_text = "\n\n".join(full_text_parts)

                # 构建初始结果
                result = ParsedDocument(
                    file_path=file_path,
                    file_name=path.name,
                    file_type="pdf",
                    total_pages=total_pages,
                    paragraphs=paragraphs,
                    raw_text=raw_text,
                )

                # 文本量充足 → 直接返回
                if not result.is_empty and not cls._is_likely_scanned(result):
                    return result

                # 文本为空或疑似扫描件 → 尝试 OCR 兜底
                ocr_result = cls._ocr_pdf(file_path, path)
                if ocr_result.parse_error:
                    # OCR 不可用，返回原始结果（可能带 parse_error 或空文本）
                    if result.is_empty:
                        return ParsedDocument(
                            file_path=file_path,
                            file_name=path.name,
                            file_type="pdf",
                            total_pages=total_pages,
                            parse_error="PDF 中未检测到文字内容，可能是扫描件（纯图片PDF），OCR 不可用: " + ocr_result.parse_error,
                        )
                    # 混合型但 OCR 不可用，返回已有的少量文本
                    result.parse_error = "PDF 文本量较少，OCR 兜底不可用: " + ocr_result.parse_error
                    return result

                if ocr_result.is_empty:
                    # OCR 也没提取到内容
                    if result.is_empty:
                        return ParsedDocument(
                            file_path=file_path,
                            file_name=path.name,
                            file_type="pdf",
                            total_pages=total_pages,
                            parse_error="PDF 中未检测到文字内容，OCR 也未识别到文本",
                        )
                    return result

                # 合并去重：PyMuPDF 已提取到足够文本的页面跳过 OCR 结果
                ocr_paragraphs_by_page: dict[int, list[dict]] = {}
                for p in ocr_result.paragraphs:
                    page = p.get("page", 0)
                    ocr_paragraphs_by_page.setdefault(page, []).append(p)

                merged_paragraphs = list(paragraphs)  # 保留 PyMuPDF 的段落
                merged_text_parts = list(full_text_parts)

                for page_key, ocr_paras in ocr_paragraphs_by_page.items():
                    # page_key 是 1-based 页码，page_text_lengths 是 0-based
                    pymupdf_chars = page_text_lengths.get(page_key - 1, 0)
                    if pymupdf_chars >= 50:
                        # 该页 PyMuPDF 已提取到足够文本，跳过 OCR 结果
                        continue
                    # 该页 PyMuPDF 文本不足，补充 OCR 结果
                    for p in ocr_paras:
                        merged_paragraphs.append(p)
                        merged_text_parts.append(p["text"])

                merged_raw_text = "\n\n".join(merged_text_parts)

                if not merged_raw_text.strip():
                    return ParsedDocument(
                        file_path=file_path,
                        file_name=path.name,
                        file_type="pdf",
                        total_pages=total_pages,
                        parse_error="PDF 中未检测到文字内容，合并 PyMuPDF 和 OCR 后仍为空",
                    )

                return ParsedDocument(
                    file_path=file_path,
                    file_name=path.name,
                    file_type="pdf",
                    total_pages=total_pages,
                    paragraphs=merged_paragraphs,
                    raw_text=merged_raw_text,
                )
            finally:
                doc.close()

        except ImportError:
            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type="pdf",
                parse_error="PyMuPDF 未安装。请执行: pip install pymupdf",
            )
        except Exception as e:
            logger.error(f"PDF 解析失败: {file_path} | {e}")
            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type="pdf",
                parse_error=f"解析失败: {e}",
            )

    @classmethod
    def _is_likely_scanned(cls, result: ParsedDocument) -> bool:
        """判断是否为扫描件：文本量与页数比值过低"""
        if result.total_pages == 0:
            return False
        chars_per_page = len(result.raw_text) / result.total_pages
        return chars_per_page < 50

    @classmethod
    def _ocr_pdf(cls, file_path: str, path: Path) -> ParsedDocument:
        """RapidOCR 识别 PDF 中的图片页面（分页流式处理，限制并发内存）"""
        try:
            from rapidocr_onnxruntime import RapidOCR
        except ImportError:
            return ParsedDocument(
                file_path=file_path, file_name=path.name, file_type="pdf",
                parse_error="RapidOCR 未安装。请执行: pip install rapidocr-onnxruntime",
            )

        try:
            import fitz  # PyMuPDF

            ocr = RapidOCR()
            doc = fitz.open(file_path)
            try:
                paragraphs = []
                full_text_parts = []

                # OCR 并发页数限制（通过环境变量配置，默认 10 页）
                _OCR_MAX_CONCURRENT = int(os.environ.get("THINKVAULT_OCR_MAX_PAGES", "10"))

                for page_num, page in enumerate(doc):
                    # 逐页渲染为图片 (200 DPI)
                    mat = fitz.Matrix(200 / 72, 200 / 72)
                    pixmap = page.get_pixmap(matrix=mat)

                    # pixmap 转为 numpy 数组供 OCR 使用
                    import numpy as np
                    img = np.frombuffer(pixmap.samples, dtype=np.uint8).reshape(
                        pixmap.height, pixmap.width, pixmap.n
                    )
                    # RapidOCR 支持 BGR/RGB，pixmap 默认 RGB
                    if pixmap.n == 4:
                        img = img[:, :, :3]  # 去掉 alpha 通道

                    # 立即释放 pixmap 内存
                    pixmap = None

                    # OCR 识别
                    ocr_result, _ = ocr(img)
                    # 释放图片内存
                    img = None

                    if ocr_result:
                        page_texts = []
                        for item in ocr_result:
                            # item: (bbox, text, confidence)
                            text = item[1].strip()
                            if text:
                                page_texts.append(text)

                        if page_texts:
                            page_content = "\n".join(page_texts)
                            paragraphs.append({
                                "page": page_num + 1,
                                "text": page_content,
                                "char_count": len(page_content),
                                "ocr": True,
                            })
                            full_text_parts.append(page_content)

                    # 每处理 _OCR_MAX_CONCURRENT 页，主动释放内存
                    if (page_num + 1) % _OCR_MAX_CONCURRENT == 0:
                        import gc
                        gc.collect()

            finally:
                total_pages = len(doc)
                doc.close()

            raw_text = "\n\n".join(full_text_parts)

            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type="pdf",
                total_pages=total_pages,
                paragraphs=paragraphs,
                raw_text=raw_text,
            )

        except Exception as e:
            logger.error(f"PDF OCR 处理失败: {file_path} | {e}")
            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type="pdf",
                parse_error=f"OCR 处理失败: {e}",
            )

    @classmethod
    def _parse_docx(cls, file_path: str, path: Path) -> ParsedDocument:
        try:
            from docx import Document

            # 带超时重试：Word 文档被其他程序锁定时，等待后可重试
            deadline = time.time() + DOCX_LOCK_TIMEOUT
            last_error = None
            while True:
                try:
                    doc = Document(file_path)
                    break
                except (IOError, PermissionError) as e:
                    last_error = e
                    if time.time() >= deadline:
                        return ParsedDocument(
                            file_path=file_path,
                            file_name=path.name,
                            file_type="docx",
                            parse_error=f"文件被占用，等待 {DOCX_LOCK_TIMEOUT}s 后仍无法打开: {e}",
                        )
                    time.sleep(0.5)
            paragraphs = []
            full_text_parts = []

            for i, para in enumerate(doc.paragraphs):
                text = para.text.strip()
                if text:
                    paragraphs.append({
                        "index": i,
                        "text": text,
                        "char_count": len(text),
                    })
                    full_text_parts.append(text)

            raw_text = "\n\n".join(full_text_parts)

            if not raw_text.strip():
                return ParsedDocument(
                    file_path=file_path,
                    file_name=path.name,
                    file_type="docx",
                    parse_error="DOCX 文档中未检测到文字内容，可能是空文档或纯图片文档",
                )

            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type="docx",
                paragraphs=paragraphs,
                raw_text=raw_text,
            )

        except ImportError:
            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type="docx",
                parse_error="python-docx 未安装。请执行: pip install python-docx",
            )
        except Exception as e:
            logger.error(f"DOCX 解析失败: {file_path} | {e}")
            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type="docx",
                parse_error=f"解析失败: {e}",
            )

    @classmethod
    def _parse_pptx(cls, file_path: str, path: Path) -> ParsedDocument:
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            paragraphs = []
            full_text_parts = []
            total_slides = len(prs.slides)

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_texts = []
                # 提取所有形状中的文本
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            text = para.text.strip()
                            if text:
                                slide_texts.append(text)
                if slide_texts:
                    slide_content = "\n".join(slide_texts)
                    paragraphs.append({
                        "page": slide_num,
                        "text": slide_content,
                        "char_count": len(slide_content),
                    })
                    full_text_parts.append(slide_content)

                # 提取表格
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        rows_data = []
                        for row in table.rows:
                            cells = [cell.text.strip() for cell in row.cells]
                            rows_data.append(" | ".join(cells))
                        if rows_data:
                            header = rows_data[0]
                            body = rows_data[1:]
                            lines = [header, " | ".join("---" for _ in header.split(" | "))]
                            lines.extend(body)
                            table_md = "\n".join(lines)
                            paragraphs.append({
                                "page": slide_num,
                                "text": table_md,
                                "char_count": len(table_md),
                                "is_table": True,
                            })
                            full_text_parts.append(table_md)

            raw_text = "\n\n".join(full_text_parts)

            if not raw_text.strip():
                return ParsedDocument(
                    file_path=file_path,
                    file_name=path.name,
                    file_type="pptx",
                    total_pages=total_slides,
                    parse_error="PPT 中未检测到文字内容，可能是纯图片幻灯片",
                )

            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type="pptx",
                total_pages=total_slides,
                paragraphs=paragraphs,
                raw_text=raw_text,
            )

        except ImportError:
            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type="pptx",
                parse_error="python-pptx 未安装。请执行: pip install python-pptx",
            )
        except Exception as e:
            logger.error(f"PPTX 解析失败: {file_path} | {e}")
            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type="pptx",
                parse_error=f"解析失败: {e}",
            )

    @classmethod
    def _parse_xlsx(cls, file_path: str, path: Path) -> ParsedDocument:
        try:
            import openpyxl

            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            paragraphs = []
            full_text_parts = []
            table_count = 0

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = list(ws.iter_rows(values_only=True))
                if not rows:
                    continue

                # 过滤全空行
                non_empty_rows = []
                for row in rows:
                    if any(cell is not None and str(cell).strip() != "" for cell in row):
                        non_empty_rows.append(row)

                if not non_empty_rows:
                    continue

                # 构建 Markdown 表格
                max_cols = max(len(row) for row in non_empty_rows)
                # 规范化每行列数
                padded_rows = []
                for row in non_empty_rows:
                    padded = list(row) + [""] * (max_cols - len(row))
                    padded_rows.append([str(c) if c is not None else "" for c in padded])

                header = padded_rows[0]
                body = padded_rows[1:]

                lines = ["## " + sheet_name, ""]
                lines.append(" | ".join(header))
                lines.append(" | ".join("---" for _ in header))
                for row in body:
                    lines.append(" | ".join(row))

                table_md = "\n".join(lines)
                table_count += 1
                paragraphs.append({
                    "page": table_count,
                    "text": table_md,
                    "char_count": len(table_md),
                    "is_table": True,
                    "sheet_name": sheet_name,
                })
                full_text_parts.append(table_md)

            wb.close()
            raw_text = "\n\n".join(full_text_parts)

            if not raw_text.strip():
                return ParsedDocument(
                    file_path=file_path,
                    file_name=path.name,
                    file_type="xlsx",
                    parse_error="Excel 文件中未检测到有效数据",
                )

            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type="xlsx",
                total_pages=table_count,
                paragraphs=paragraphs,
                raw_text=raw_text,
            )

        except ImportError:
            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type="xlsx",
                parse_error="openpyxl 未安装。请执行: pip install openpyxl",
            )
        except Exception as e:
            logger.error(f"XLSX 解析失败: {file_path} | {e}")
            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type="xlsx",
                parse_error=f"解析失败: {e}",
            )

    @classmethod
    def _parse_txt(cls, file_path: str, path: Path) -> ParsedDocument:
        # 大文件保护：超过 50MB 拒绝一次性读取
        _MAX_TXT_SIZE = 50 * 1024 * 1024  # 50MB
        if path.stat().st_size > _MAX_TXT_SIZE:
            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type="txt",
                parse_error=f"文件过大 ({path.stat().st_size // (1024*1024)}MB)，最大支持 50MB",
            )
        try:
            text = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            try:
                text = path.read_text(encoding="gbk")
            except Exception:
                return ParsedDocument(
                    file_path=file_path,
                    file_name=path.name,
                    file_type="txt",
                    parse_error="无法识别文件编码，请转换为 UTF-8",
                )

        paragraphs = [
            {"index": i, "text": p.strip(), "char_count": len(p.strip())}
            for i, p in enumerate(text.split("\n\n"))
            if p.strip()
        ]

        return ParsedDocument(
            file_path=file_path,
            file_name=path.name,
            file_type="txt",
            paragraphs=paragraphs,
            raw_text=text,
        )

    @classmethod
    def _parse_markdown(cls, file_path: str, path: Path) -> ParsedDocument:
        doc = cls._parse_txt(file_path, path)
        doc.file_type = "md"
        return doc

    # ── 音频解析 ──────────────────────────────────────────────────

    @classmethod
    def _parse_audio(cls, file_path: str, path: Path) -> ParsedDocument:
        """使用 faster-whisper 转录音频文件，按 30 秒分段"""
        suffix = path.suffix.lower()
        file_type = suffix.lstrip(".")

        if not _whisper_available:
            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type=file_type,
                parse_error="faster-whisper 未安装。请执行: pip install 'thinkvault[media]' 或 pip install faster-whisper",
            )

        try:
            # 懒加载 Whisper 模型（首次调用时加载，后续复用）
            model_size = os.environ.get("THINKVAULT_WHISPER_MODEL", "base")
            device = os.environ.get("THINKVAULT_WHISPER_DEVICE", "cpu")
            compute_type = os.environ.get("THINKVAULT_WHISPER_COMPUTE_TYPE", "int8")

            # 缓存 Whisper 模型实例，避免每次调用重新加载（5-15秒）
            cache_key = f"{model_size}_{device}_{compute_type}"
            model = _whisper_model_cache.get(cache_key)
            if model is None:
                with _whisper_lock:
                    model = _whisper_model_cache.get(cache_key)
                    if model is None:
                        model = WhisperModel(model_size, device=device, compute_type=compute_type)
                        _whisper_model_cache[cache_key] = model

            segments_iter, info = model.transcribe(file_path, beam_size=5)

            paragraphs = []
            full_text_parts = []

            for segment in segments_iter:
                start_sec = segment.start
                end_sec = segment.end
                text = segment.text.strip()

                if text:
                    # 格式化为 [MM:SS - MM:SS] 文本
                    start_min, start_s = divmod(int(start_sec), 60)
                    end_min, end_s = divmod(int(end_sec), 60)
                    timestamp = f"[{start_min:02d}:{start_s:02d} - {end_min:02d}:{end_s:02d}]"
                    segment_text = f"{timestamp} {text}"

                    paragraphs.append({
                        "page": 1,
                        "text": segment_text,
                        "char_count": len(segment_text),
                        "start": start_sec,
                        "end": end_sec,
                    })
                    full_text_parts.append(segment_text)

            raw_text = "\n\n".join(full_text_parts)

            if not raw_text.strip():
                return ParsedDocument(
                    file_path=file_path,
                    file_name=path.name,
                    file_type=file_type,
                    parse_error="音频转录结果为空，可能是无声文件或无法识别的语音",
                )

            # 音频时长转换为"页数"（每30秒算1页，向上取整）
            total_duration = info.duration if hasattr(info, "duration") else 0
            total_pages = max(1, int(total_duration / 30) + (1 if total_duration % 30 > 0 else 0))

            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type=file_type,
                total_pages=total_pages,
                paragraphs=paragraphs,
                raw_text=raw_text,
            )

        except Exception as e:
            logger.error(f"音频解析失败: {file_path} | {e}")
            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type=file_type,
                parse_error=f"音频转录失败: {e}",
            )

    # ── 视频解析 ──────────────────────────────────────────────────

    @classmethod
    def _parse_video(cls, file_path: str, path: Path) -> ParsedDocument:
        """先用 ffmpeg 提取音频轨道，再调用 _parse_audio 转录"""
        suffix = path.suffix.lower()
        file_type = suffix.lstrip(".")

        if not _ffmpeg_available:
            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type=file_type,
                parse_error="ffmpeg 未安装或不在 PATH 中。请安装 ffmpeg 后重试",
            )

        if not _whisper_available:
            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type=file_type,
                parse_error="faster-whisper 未安装。请执行: pip install 'thinkvault[media]' 或 pip install faster-whisper",
            )

        temp_audio_path = None
        try:
            # 用 ffmpeg 提取音频到临时 WAV 文件
            temp_dir = tempfile.mkdtemp(prefix="thinkvault_video_")
            temp_audio_path = os.path.join(temp_dir, "audio_extract.wav")

            import subprocess
            cmd = [
                "ffmpeg", "-i", file_path,
                "-vn",              # 不要视频
                "-acodec", "pcm_s16le",  # 16-bit PCM
                "-ar", "16000",     # 16kHz 采样率（Whisper 最优）
                "-ac", "1",         # 单声道
                "-y",               # 覆盖输出
                temp_audio_path,
            ]
            result = subprocess.run(
                cmd,
                capture_output=True,
                text=True,
                timeout=300,  # 5 分钟超时
            )
            if result.returncode != 0:
                return ParsedDocument(
                    file_path=file_path,
                    file_name=path.name,
                    file_type=file_type,
                    parse_error=f"ffmpeg 提取音频失败: {result.stderr[:500]}",
                )

            # 调用音频转录
            audio_doc = cls._parse_audio(temp_audio_path, Path(temp_audio_path))
            # 替换文件名和路径为原始视频文件信息
            audio_doc.file_path = file_path
            audio_doc.file_name = path.name
            audio_doc.file_type = file_type

            if audio_doc.parse_error:
                audio_doc.parse_error = f"视频音频转录失败: {audio_doc.parse_error}"

            return audio_doc

        except subprocess.TimeoutExpired:
            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type=file_type,
                parse_error="ffmpeg 提取音频超时（5分钟），文件可能过大",
            )
        except Exception as e:
            logger.error(f"视频解析失败: {file_path} | {e}")
            return ParsedDocument(
                file_path=file_path,
                file_name=path.name,
                file_type=file_type,
                parse_error=f"视频解析失败: {e}",
            )
        finally:
            # 清理临时文件
            if temp_audio_path and os.path.exists(temp_audio_path):
                try:
                    temp_dir = os.path.dirname(temp_audio_path)
                    shutil.rmtree(temp_dir, ignore_errors=True)
                except Exception:
                    pass
